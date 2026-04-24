import streamlit as st
import io
import base64
import os
import json
import warnings
import pypdf
import pandas as pd
import datetime
import uuid

from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import HumanMessage, AIMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter

warnings.filterwarnings("ignore")

st.set_page_config(page_title="StudyMate AI", page_icon="🎓", layout="wide")

# ================================================================
# PERSISTENCE LAYER
# Sessions stored as JSON in ./sessions/
# Schema: { id, name, created_at, updated_at,
#           input_mode, topic, messages: [{role, content}] }
#
# Save policy: sessions are written to disk only when there is
# content worth keeping — not on every rerun or UI interaction.
# ================================================================
SESSIONS_DIR = "sessions"
os.makedirs(SESSIONS_DIR, exist_ok=True)

MAX_CHAT_HISTORY = 20  # max LangChain message objects kept in memory


def _session_path(session_id: str) -> str:
    return os.path.join(SESSIONS_DIR, f"{session_id}.json")


def list_sessions() -> list[dict]:
    """Return all saved sessions sorted by most recently updated."""
    sessions = []
    for fname in os.listdir(SESSIONS_DIR):
        if fname.endswith(".json"):
            try:
                with open(os.path.join(SESSIONS_DIR, fname), "r") as f:
                    sessions.append(json.load(f))
            except Exception:
                pass
    return sorted(sessions, key=lambda s: s.get("updated_at", ""), reverse=True)


def save_session(session: dict) -> None:
    """Write a session dict to disk, updating updated_at."""
    session["updated_at"] = datetime.datetime.now().isoformat()
    with open(_session_path(session["id"]), "w") as f:
        json.dump(session, f, indent=2)


def load_session(session_id: str) -> dict | None:
    path = _session_path(session_id)
    if os.path.exists(path):
        with open(path, "r") as f:
            return json.load(f)
    return None


def delete_session(session_id: str) -> None:
    path = _session_path(session_id)
    if os.path.exists(path):
        os.remove(path)


def new_session(name: str = "", input_mode: str = "", topic: str = "") -> dict:
    """Create a fresh session dict. Not saved to disk until persist_message() is called."""
    now = datetime.datetime.now().isoformat()
    return {
        "id": str(uuid.uuid4()),
        "name": name or f"Session {datetime.datetime.now().strftime('%b %d, %H:%M')}",
        "created_at": now,
        "updated_at": now,
        "input_mode": input_mode,
        "topic": topic,
        "messages": [],
    }


def export_session_as_txt(session: dict) -> str:
    lines = [
        f"StudyMate AI — {session['name']}",
        f"Created: {session['created_at'][:16].replace('T', ' ')}",
        f"Topic / Source: {session.get('topic', 'N/A')}",
        "=" * 60,
        "",
    ]
    for msg in session.get("messages", []):
        role = "You" if msg["role"] == "user" else "StudyMate AI"
        lines.append(f"[{role}]")
        lines.append(msg["content"])
        lines.append("")
    return "\n".join(lines)


# ================================================================
# APP CONFIG / CONSTANTS
# ================================================================
try:
    groq_api_key = st.secrets["GROQ_API_KEY"]
except Exception:
    st.error("⚠️ Service unavailable. Please contact the administrator.")
    st.stop()

PERSONAS = {
    "Summarize": "You summarize content into concise bullet points. Build on previous summaries if the student asks follow-up questions.",
    "Professor": "You are a friendly college professor who explains concepts using analogies and simple examples. Remember what you have already taught the student and build on it.",
    "Study Guide": "You extract key definitions and important concepts to create structured study guides. Remember previously covered topics and expand on them when asked.",
    "Socratic Tutor": "You are a Socratic tutor. Do NOT give final answers. Guide the student using hints and thoughtful questions. Remember the student's previous responses and build your next question on them.",
    "Quiz Master": (
        "You are an engaging, encouraging, and sharp-witted quiz master. Your goal is to test knowledge, reinforce learning, and keep the user motivated. "
        "You ask clear, concise questions based on previous topics, provide immediate, constructive feedback on answers, and keep score if requested. "
        "You tailor the difficulty level to the student's progress and explain why an answer is correct or incorrect to deepen understanding."
    ),
    "ELI5 Simplifier": (
        "You are an expert at explaining complex topics in the simplest way possible — as if explaining to a curious 10-year-old. "
        "Use everyday analogies, short sentences, and familiar examples. Avoid jargon entirely; if a technical term is unavoidable, immediately explain it in plain language. "
        "Break everything down into small, digestible pieces. Use comparisons to things from daily life (food, games, sports, toys) to make abstract ideas concrete. "
        "Keep a warm, encouraging, and enthusiastic tone — learning should feel easy and fun, never overwhelming. "
        "If the student seems confused, try a completely different analogy. Never say 'it's complicated' — everything can be made simple."
    ),
}

SUPPORTED_EXTENSIONS = [
    "pdf", "docx", "doc", "odt", "rtf", "pptx", "ppt",
    "xlsx", "xls", "xlsm", "ods", "csv", "tsv",
    "txt", "md", "py", "js", "ts", "html", "css", "java",
    "c", "cpp", "cs", "go", "rs", "rb", "php", "swift",
    "json", "jsonl", "xml", "yaml", "yml", "toml", "ini", "log", "ipynb",
]

FILE_TYPE_LABELS = {
    "pdf": "📄 PDF", "docx": "📝 Word", "doc": "📝 Word (legacy)",
    "odt": "📝 OpenDocument Text", "rtf": "📝 Rich Text",
    "pptx": "📊 PowerPoint", "ppt": "📊 PowerPoint (legacy)",
    "xlsx": "📈 Excel", "xls": "📈 Excel (legacy)",
    "xlsm": "📈 Excel (macro)", "ods": "📈 OpenDocument Sheet",
    "csv": "📋 CSV", "tsv": "📋 TSV",
    "txt": "📃 Text", "md": "📃 Markdown", "ipynb": "📓 Jupyter Notebook",
    "json": "🔧 JSON", "jsonl": "🔧 JSONL", "xml": "🔧 XML",
    "yaml": "🔧 YAML", "yml": "🔧 YAML", "toml": "🔧 TOML",
    "ini": "🔧 Config", "log": "📋 Log",
    "py": "🐍 Python", "js": "⚡ JavaScript", "ts": "⚡ TypeScript",
    "html": "🌐 HTML", "css": "🎨 CSS", "java": "☕ Java",
    "c": "⚙️ C", "cpp": "⚙️ C++", "cs": "⚙️ C#",
    "go": "🐹 Go", "rs": "🦀 Rust", "rb": "💎 Ruby",
    "php": "🐘 PHP", "swift": "🍎 Swift",
}


# ================================================================
# FILE EXTRACTION
# ================================================================
def extract_text_from_file(uploaded_file) -> str:
    filename = uploaded_file.name
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    raw = uploaded_file.read()
    buf = io.BytesIO(raw)

    if ext == "pdf":
        reader = pypdf.PdfReader(buf)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {i+1}]\n{text.strip()}")
        if not pages:
            raise ValueError("Could not extract text from this PDF. It may be scanned or image-based.")
        return "\n\n".join(pages)

    elif ext in ("docx", "odt", "rtf", "epub"):
        tmp = f"/tmp/_studymate_upload.{ext}"
        with open(tmp, "wb") as f:
            f.write(raw)
        try:
            import docx2txt
            text = docx2txt.process(tmp)
            if text and text.strip():
                return text.strip()
        except Exception:
            pass
        try:
            import subprocess
            result = subprocess.run(["pandoc", tmp, "-t", "plain"], capture_output=True, text=True, timeout=30)
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except Exception:
            pass
        raise ValueError(f"Could not extract text from {ext.upper()} file.")

    elif ext in ("pptx", "ppt"):
        from pptx import Presentation
        prs = Presentation(buf)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        if not slides:
            raise ValueError("No text found in presentation.")
        return "\n\n".join(slides)

    elif ext in ("xlsx", "xlsm"):
        df_dict = pd.read_excel(buf, sheet_name=None, engine="openpyxl")
        return "\n\n".join(f"[Sheet: {n}]\n{df.to_string(index=False)}" for n, df in df_dict.items())

    elif ext == "ods":
        df_dict = pd.read_excel(buf, sheet_name=None, engine="odf")
        return "\n\n".join(f"[Sheet: {n}]\n{df.to_string(index=False)}" for n, df in df_dict.items())

    elif ext == "xls":
        df_dict = pd.read_excel(buf, sheet_name=None, engine="xlrd")
        return "\n\n".join(f"[Sheet: {n}]\n{df.to_string(index=False)}" for n, df in df_dict.items())

    elif ext in ("csv", "tsv"):
        df = pd.read_csv(buf, sep="\t" if ext == "tsv" else ",")
        return df.to_string(index=False)

    elif ext == "json":
        return json.dumps(json.load(buf), indent=2)

    elif ext == "jsonl":
        lines = raw.decode("utf-8", errors="replace").strip().split("\n")
        parsed = []
        for line in lines[:500]:
            try:
                parsed.append(json.dumps(json.loads(line)))
            except Exception:
                parsed.append(line)
        return "\n".join(parsed)

    elif ext == "ipynb":
        nb = json.loads(raw.decode("utf-8", errors="replace"))
        cells = []
        for cell in nb.get("cells", []):
            src = "".join(cell.get("source", []))
            if src.strip():
                label = "# CODE CELL" if cell.get("cell_type") == "code" else "# MARKDOWN"
                cells.append(f"{label}\n{src.strip()}")
        return "\n\n".join(cells)

    elif ext in ("txt", "md", "py", "js", "ts", "html", "css", "java",
                 "c", "cpp", "cs", "go", "rs", "rb", "php", "swift",
                 "xml", "yaml", "yml", "toml", "ini", "log"):
        return raw.decode("utf-8", errors="replace")

    else:
        decoded = raw.decode("utf-8", errors="replace")
        if decoded.strip():
            return decoded
        raise ValueError(f"Unsupported or unreadable file type: .{ext}")


def build_vectorstore(text: str, source_name: str):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=400, chunk_overlap=50)
    docs = text_splitter.create_documents([text], metadatas=[{"source": source_name}])
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    return FAISS.from_documents(docs, embeddings)


def run_rag_chat(user_query: str, retriever, llm, persona: str) -> str:
    prompt = ChatPromptTemplate.from_messages([
        ("system", f"""{persona}

You have access to the following document context to answer questions.
Always use this context to ground your answers.
If the answer is not in the context, say "The answer is not found in the document."

Document Context:
{{context}}"""),
        MessagesPlaceholder(variable_name="chat_history"),
        ("human", "{question}")
    ])
    retrieved_docs = retriever.invoke(user_query)
    context = "\n\n".join(doc.page_content for doc in retrieved_docs)
    chain = prompt | llm | StrOutputParser()
    return chain.invoke({
        "context": context,
        "chat_history": st.session_state.chat_history,
        "question": user_query,
    })


def lc_history_from_messages(messages: list[dict]) -> list:
    """Rebuild LangChain message objects from saved message dicts."""
    history = []
    for m in messages:
        if m["role"] == "user":
            history.append(HumanMessage(content=m["content"]))
        elif m["role"] == "assistant":
            history.append(AIMessage(content=m["content"]))
    return history


# ================================================================
# SESSION STATE BOOTSTRAP
# current_session: dict | None  — None means "not started yet"
# messages:        list[dict]   — independent copy, never shared ref
# chat_history:    list[LC msg] — trimmed to MAX_CHAT_HISTORY pairs
# ================================================================
if "current_session" not in st.session_state:
    st.session_state.current_session = None
if "messages" not in st.session_state:
    st.session_state.messages = []
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []


def _ensure_session() -> None:
    """Lazily create a session the first time content is about to be saved.
    This prevents ghost sessions for users who open the app but never chat."""
    if st.session_state.current_session is None:
        sess = new_session()
        st.session_state.current_session = sess
        # Don't save yet — persist_message() will save on first message.


def persist_message(role: str, content: str) -> None:
    """Append a message to both in-memory lists and save the session to disk.

    Save policy:
    - Ensures a session exists (lazy creation).
    - Saves to disk once per message — not on every UI interaction.
    - Uses a fresh copy of messages so session dict and st.session_state.messages
      are always independent objects (no shared-reference mutations).
    """
    _ensure_session()

    msg = {"role": role, "content": content}

    # Keep UI list and session list as separate objects — no shared references.
    st.session_state.messages.append(msg)
    st.session_state.current_session["messages"] = list(st.session_state.messages)

    # Trim LangChain history to avoid unbounded token growth.
    st.session_state.chat_history = st.session_state.chat_history[-MAX_CHAT_HISTORY:]

    save_session(st.session_state.current_session)


def _load_session_into_state(sess: dict) -> None:
    """Switch the active session. Uses list() to guarantee independent copies."""
    st.session_state.current_session = dict(sess)  # top-level copy
    st.session_state.messages = list(sess.get("messages", []))  # independent list
    st.session_state.chat_history = lc_history_from_messages(st.session_state.messages)[-MAX_CHAT_HISTORY:]
    # Clear stale file / image state — user must re-upload
    for k in ["vectorstore", "file_key", "image_data", "image_media_type", "image_hash"]:
        st.session_state.pop(k, None)
    # Only restore free_topic when the session was actually in topic mode
    if sess.get("input_mode") == "✏️ Enter a Topic":
        st.session_state.free_topic = sess.get("topic", "")
    else:
        st.session_state.pop("free_topic", None)


# ================================================================
# SIDEBAR — Sessions Manager + Settings
# ================================================================
st.sidebar.title("🎓 StudyMate AI")

# ── Settings ─────────────────────────────────────────────────
st.sidebar.header("⚙️ Settings")
mode = st.sidebar.radio("Study Mode:", list(PERSONAS.keys()))
temperature = st.sidebar.slider("Creativity Level", 0.0, 1.0, 0.3)
model_choice = st.sidebar.selectbox("Model:", [
    "llama-3.3-70b-versatile",
    "llama-3.1-8b-instant",
    "gemma2-9b-it",
    "deepseek-r1-distill-llama-70b",
])

st.sidebar.markdown("---")

# ── Session controls ─────────────────────────────────────────
st.sidebar.header("💾 Sessions")

col_new, col_clear = st.sidebar.columns(2)
with col_new:
    if st.button("＋ New", use_container_width=True, help="Start a brand-new session"):
        # Reset to a blank slate. A new session dict is created lazily on first message.
        st.session_state.current_session = None
        st.session_state.messages = []
        st.session_state.chat_history = []
        for k in ["vectorstore", "file_key", "image_data", "image_media_type", "image_hash", "free_topic"]:
            st.session_state.pop(k, None)
        st.rerun()

with col_clear:
    if st.button("🗑️ Clear", use_container_width=True, help="Clear current chat (keeps session)"):
        st.session_state.messages = []
        st.session_state.chat_history = []
        if st.session_state.current_session:
            st.session_state.current_session["messages"] = []
            save_session(st.session_state.current_session)
        st.rerun()

# ── Rename current session ────────────────────────────────────
if st.session_state.current_session:
    with st.sidebar.expander("✏️ Rename current session"):
        new_name = st.text_input(
            "Session name",
            value=st.session_state.current_session.get("name", ""),
            key="rename_input",
            label_visibility="collapsed",
        )
        if st.button("Save name", key="save_name"):
            trimmed = new_name.strip()
            if trimmed:
                st.session_state.current_session["name"] = trimmed
                save_session(st.session_state.current_session)
            st.rerun()

# ── Download current session ──────────────────────────────────
if st.session_state.current_session and st.session_state.messages:
    txt = export_session_as_txt(st.session_state.current_session)
    safe_name = st.session_state.current_session["name"].replace(" ", "_")[:40]
    st.sidebar.download_button(
        label="⬇️ Download notes (.txt)",
        data=txt,
        file_name=f"{safe_name}.txt",
        mime="text/plain",
        use_container_width=True,
    )

st.sidebar.markdown("---")

# ── Saved sessions list ───────────────────────────────────────
st.sidebar.subheader("📚 Saved Sessions")
all_sessions = list_sessions()

if not all_sessions:
    st.sidebar.caption("No saved sessions yet. Start chatting to auto-save!")
else:
    for sess in all_sessions:
        is_active = (
            st.session_state.current_session is not None
            and st.session_state.current_session["id"] == sess["id"]
        )
        label = ("▶ " if is_active else "") + sess["name"]
        updated = sess.get("updated_at", "")[:16].replace("T", " ")
        msg_count = len(sess.get("messages", []))

        with st.sidebar.container():
            col_btn, col_del = st.sidebar.columns([4, 1])
            with col_btn:
                if st.button(
                    label,
                    key=f"load_{sess['id']}",
                    use_container_width=True,
                    type="primary" if is_active else "secondary",
                    help=f"{updated} · {msg_count} msgs",
                ):
                    if not is_active:
                        _load_session_into_state(sess)
                        st.rerun()
            with col_del:
                if st.button("🗑", key=f"del_{sess['id']}", help="Delete this session"):
                    delete_session(sess["id"])
                    if is_active:
                        st.session_state.current_session = None
                        st.session_state.messages = []
                        st.session_state.chat_history = []
                    st.rerun()

with st.sidebar.expander("📂 Supported File Types"):
    st.markdown("""
**Documents:** PDF, Word, ODT, RTF  
**Presentations:** PowerPoint  
**Spreadsheets:** Excel, ODS, CSV, TSV  
**Code & Text:** Python, JS, TS, HTML, CSS, Java, C/C++, Go, Rust, Ruby, PHP, Swift, MD, TXT  
**Data:** JSON, JSONL, XML, YAML, TOML, INI, LOG  
**Notebooks:** Jupyter (.ipynb)
""")


# ================================================================
# MAIN AREA
# ================================================================
st.title("🎓 StudyMate AI")
st.markdown("### Your Personalized AI Study Companion")

# Session label — shows "New session" until first message is saved
if st.session_state.current_session:
    st.caption(f"📝 Session: **{st.session_state.current_session['name']}**")
else:
    st.caption("📝 New session — start chatting to save it automatically.")

st.markdown("---")
input_mode = st.radio(
    "📚 How would you like to study today?",
    ["📁 Upload a File", "🖼️ Upload Image", "📷 Take a Photo", "✏️ Enter a Topic"],
    horizontal=True,
)

# Persist input_mode change only when there's an active saved session
if (
    st.session_state.current_session is not None
    and st.session_state.current_session.get("input_mode") != input_mode
    and st.session_state.current_session.get("messages")  # don't save empty-session metadata churn
):
    st.session_state.current_session["input_mode"] = input_mode
    save_session(st.session_state.current_session)


# ==============================
# MODE 1: Universal File Upload
# ==============================
if input_mode == "📁 Upload a File":
    uploaded_file = st.file_uploader(
        "Upload any document, spreadsheet, code file, or data file",
        type=SUPPORTED_EXTENSIONS,
        label_visibility="collapsed",
    )

    if uploaded_file:
        ext = uploaded_file.name.rsplit(".", 1)[-1].lower() if "." in uploaded_file.name else "?"
        file_label = FILE_TYPE_LABELS.get(ext, f"📄 .{ext.upper()}")
        file_key = f"{uploaded_file.name}_{uploaded_file.size}"

        if "vectorstore" not in st.session_state or st.session_state.get("file_key") != file_key:
            with st.spinner(f"🔄 Processing your {file_label} file..."):
                try:
                    raw_text = extract_text_from_file(uploaded_file)
                except ValueError as e:
                    st.error(f"❌ {e}")
                    st.stop()
                except Exception as e:
                    st.error(f"❌ Unexpected error reading file: {e}")
                    st.stop()

                if not raw_text or not raw_text.strip():
                    st.error("❌ No readable text could be extracted from this file.")
                    st.stop()

                vectorstore = build_vectorstore(raw_text, uploaded_file.name)
                st.session_state.vectorstore = vectorstore
                st.session_state.file_key = file_key

            word_count = len(raw_text.split())
            st.success(f"✅ {file_label} ready! (~{word_count:,} words extracted). Ask your first question below.")

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    if uploaded_file and "vectorstore" in st.session_state:
        user_query = st.chat_input("❓ Ask a question about your file...")
        if user_query:
            with st.chat_message("user"):
                st.markdown(user_query)
            persist_message("user", user_query)

            # Tag the session with the source file (once per file, not every message)
            if st.session_state.current_session.get("topic") != uploaded_file.name:
                st.session_state.current_session["topic"] = uploaded_file.name
                save_session(st.session_state.current_session)

            llm = ChatGroq(model=model_choice, temperature=temperature, api_key=groq_api_key)
            retriever = st.session_state.vectorstore.as_retriever(search_kwargs={"k": 3})

            with st.chat_message("assistant"):
                with st.spinner("🧠 Thinking..."):
                    response = run_rag_chat(user_query, retriever, llm, PERSONAS[mode])
                st.markdown(response)

            persist_message("assistant", response)
            st.session_state.chat_history.append(HumanMessage(content=user_query))
            st.session_state.chat_history.append(AIMessage(content=response))
    elif not st.session_state.messages:
        st.info("📁 Upload any file to get started. See **Supported File Types** in the sidebar.")


# ==============================
# MODE 2 & 3: Image / Camera
# ==============================
elif input_mode in ["🖼️ Upload Image", "📷 Take a Photo"]:
    image_data = None
    media_type = "image/jpeg"

    if input_mode == "🖼️ Upload Image":
        uploaded_img = st.file_uploader(
            "Upload an image (notes, diagrams, textbook pages, etc.)",
            type=["jpg", "jpeg", "png", "webp"],
            label_visibility="collapsed",
        )
        if uploaded_img:
            image_data = uploaded_img.read()
            media_type = uploaded_img.type or "image/jpeg"
            st.image(image_data, caption="Uploaded Image", use_container_width=True)
    else:
        camera_image = st.camera_input("📷 Point your camera at anything you want to learn about!")
        if camera_image:
            image_data = camera_image.read()
            media_type = "image/jpeg"

    if image_data:
        img_hash = hash(image_data)
        if st.session_state.get("image_hash") != img_hash:
            st.session_state.image_data = base64.standard_b64encode(image_data).decode("utf-8")
            st.session_state.image_media_type = media_type
            st.session_state.image_hash = img_hash
            # Reset conversation when a new image is loaded
            st.session_state.messages = []
            st.session_state.chat_history = []
            st.success("✅ Image ready! Ask anything about what's in the image.")

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    if st.session_state.get("image_data"):
        user_query = st.chat_input("❓ Ask a question about the image...")
        if user_query:
            with st.chat_message("user"):
                st.markdown(user_query)
            persist_message("user", user_query)

            # Tag the session with image topic (once)
            if st.session_state.current_session.get("topic") != "Image upload":
                st.session_state.current_session["topic"] = "Image upload"
                save_session(st.session_state.current_session)

            system_prompt = (
                f"{PERSONAS[mode]}\n\n"
                "The student has shared an image. Analyze it thoroughly and answer questions about it. "
                "Use the visual content (text, diagrams, equations, labels, etc.) as your primary reference."
            )

            # Always include the image as the first user turn so the model retains
            # visual context across the full conversation, not just on turn 1.
            image_turn = {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{st.session_state.image_media_type};base64,{st.session_state.image_data}"
                        },
                    },
                    {"type": "text", "text": "Please analyse this image. I will ask you questions about it."},
                ],
            }

            messages_to_send = [{"role": "system", "content": system_prompt}, image_turn]

            # Replay prior turns (text-only — image is already injected above)
            for msg in st.session_state.chat_history:
                role = "user" if isinstance(msg, HumanMessage) else "assistant"
                messages_to_send.append({"role": role, "content": msg.content})

            messages_to_send.append({"role": "user", "content": user_query})

            with st.chat_message("assistant"):
                with st.spinner("🧠 Analyzing image..."):
                    from groq import Groq
                    client = Groq(api_key=groq_api_key)
                    completion = client.chat.completions.create(
                        model="meta-llama/llama-4-scout-17b-16e-instruct",
                        messages=messages_to_send,
                        temperature=temperature,
                        max_tokens=1024,
                    )
                    response = completion.choices[0].message.content
                st.markdown(response)

            persist_message("assistant", response)
            st.session_state.chat_history.append(HumanMessage(content=user_query))
            st.session_state.chat_history.append(AIMessage(content=response))
    elif not st.session_state.messages:
        if input_mode == "🖼️ Upload Image":
            st.info("🖼️ Upload an image of your notes, a diagram, or a textbook page to get started.")
        else:
            st.info("📷 Use the camera above to capture anything you want to learn about.")


# ==============================
# MODE 4: Free Topic
# ==============================
elif input_mode == "✏️ Enter a Topic":
    st.markdown("#### 💡 What do you want to learn about today?")

    col1, col2 = st.columns([3, 1])
    with col1:
        topic_input = st.text_input(
            "Enter any topic, concept, or question:",
            placeholder="e.g. Photosynthesis, The French Revolution, Python decorators...",
            label_visibility="collapsed",
        )
    with col2:
        set_topic = st.button("🚀 Start Studying", use_container_width=True)

    if set_topic and topic_input.strip():
        topic = topic_input.strip()
        st.session_state.free_topic = topic
        st.session_state.messages = []
        st.session_state.chat_history = []
        for k in ["vectorstore", "file_key", "image_data", "image_media_type", "image_hash"]:
            st.session_state.pop(k, None)

        # Start a fresh named session for this topic.
        # Don't save to disk yet — persist_message() will do that on first message.
        sess = new_session(name=topic[:50], input_mode=input_mode, topic=topic)
        st.session_state.current_session = sess

        st.rerun()

    current_topic = st.session_state.get("free_topic", "")

    if current_topic:
        st.success(f"📖 Currently studying: **{current_topic}**")

        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

        # Auto-send a starter prompt only on the very first turn
        default_prompt = f"Let's start! Help me learn about: {current_topic}" if not st.session_state.messages else None
        user_query = st.chat_input("❓ Ask anything about this topic...") or (
            default_prompt if not st.session_state.messages else None
        )

        if user_query:
            llm = ChatGroq(model=model_choice, temperature=temperature, api_key=groq_api_key)
            prompt = ChatPromptTemplate.from_messages([
                ("system", f"""{PERSONAS[mode]}

The student wants to learn about: {current_topic}

Use your broad knowledge to teach this topic thoroughly.
If the student asks something off-topic, gently guide them back to {current_topic} unless it's a natural extension."""),
                MessagesPlaceholder(variable_name="chat_history"),
                ("human", "{question}"),
            ])
            chain = prompt | llm | StrOutputParser()

            # Only show the user bubble for real queries, not the silent auto-starter
            if user_query != default_prompt:
                with st.chat_message("user"):
                    st.markdown(user_query)

            with st.chat_message("assistant"):
                with st.spinner("🧠 Thinking..."):
                    response = chain.invoke({
                        "chat_history": st.session_state.chat_history,
                        "question": user_query,
                    })
                st.markdown(response)

            # Persist both turns (user first, even for the hidden starter)
            persist_message("user", user_query)
            persist_message("assistant", response)
            st.session_state.chat_history.append(HumanMessage(content=user_query))
            st.session_state.chat_history.append(AIMessage(content=response))
            st.rerun()
    else:
        st.info("✏️ Type a topic above and click **Start Studying** to begin — no files needed!")